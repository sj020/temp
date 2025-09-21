# pip install python-pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def replace_in_paragraph(paragraph, placeholder, replacement):
    """
    Replace all occurrences of placeholder in the paragraph, handling cases
    where the placeholder spans multiple runs.
    Replacement inherits formatting of the run where the placeholder starts.
    """
    # keep trying until no placeholder remains in this paragraph
    while True:
        runs = list(paragraph.runs)
        if not runs:
            return
        full = ''.join([r.text or '' for r in runs])
        idx = full.find(placeholder)
        if idx == -1:
            return

        plen = len(placeholder)
        # compute cumulative lengths of runs
        cum = []
        total = 0
        for r in runs:
            txt = r.text or ''
            total += len(txt)
            cum.append(total)

        start = idx
        end = idx + plen  # exclusive

        # find start run index
        for si, c in enumerate(cum):
            if c > start:
                start_i = si
                break
        # find end run index
        for ei, c in enumerate(cum):
            if c >= end:
                end_i = ei
                break

        start_run = runs[start_i]
        end_run = runs[end_i]

        # offsets inside start/end runs
        start_offset = start - (cum[start_i - 1] if start_i > 0 else 0)
        end_offset = end - (cum[end_i - 1] if end_i > 0 else 0)

        if start_i == end_i:
            # placeholder entirely inside one run -> simple replacement (keeps that run's formatting)
            s = start_run.text or ''
            start_run.text = s[:start_offset] + replacement + s[end_offset:]
        else:
            # spanning runs:
            start_txt = start_run.text or ''
            end_txt = end_run.text or ''

            prefix = start_txt[:start_offset]
            suffix = end_txt[end_offset:]

            # put replacement into the start_run (inherits its formatting)
            start_run.text = prefix + replacement
            # clear intermediate runs
            for i in range(start_i + 1, end_i):
                runs[i].text = ''
            # set suffix into end_run
            end_run.text = suffix

def replace_in_shape(shape, placeholder, replacement):
    """
    Recursively replace text in a shape (text frames, tables, groups).
    """
    # Text in regular text-carrying shapes
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            replace_in_paragraph(para, placeholder, replacement)

    # Tables (graphic frames)
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    replace_in_paragraph(para, placeholder, replacement)

    # Group shapes: recurse into contained shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shp in shape.shapes:
            replace_in_shape(shp, placeholder, replacement)

def replace_placeholder_in_presentation(input_path, output_path, placeholder, replacement):
    prs = Presentation(input_path)

    # Replace inside slide layouts & master (in case placeholder sits in layout/master)
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            replace_in_shape(shape, placeholder, replacement)
    for shape in prs.slide_master.shapes:
        replace_in_shape(shape, placeholder, replacement)

    # Replace on slides and notes
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, placeholder, replacement)

        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                replace_in_shape(shape, placeholder, replacement)

    prs.save(output_path)
    print(f"Saved replaced PPTX -> {output_path}")

import os
from pdf2image import convert_from_path, convert_from_bytes

def pdf_to_images(pdf_path, output_folder, fmt='PNG', dpi=200, poppler_path=None):
    """
    Convert each page of pdf_path into an image in output_folder.

    Arguments:
        pdf_path (str): path to input PDF file.
        output_folder (str): directory where page images will be saved.
        fmt (str): image format, e.g. 'PNG', 'JPEG'.
        dpi (int): resolution in DPI for rendering.
        poppler_path (str or None): path to poppler bin dir on systems where required (Windows).
    """

    os.makedirs(output_folder, exist_ok=True)

    if poppler_path:
        pages = convert_from_path(pdf_path, dpi=dpi, fmt=fmt, poppler_path=poppler_path)
    else:
        pages = convert_from_path(pdf_path, dpi=dpi, fmt=fmt)

    for i, page in enumerate(pages, start=1):
        fname = os.path.join(output_folder, f"page_{i}.{ fmt.lower() }")
        page.save(fname, fmt)
        print(f"Saved {fname}")

if __name__ == "__main__":
    pdf_file = "filled.pdf"         # your PDF file
    out_dir = "pdf_pages_images"     # folder to store images
    # If you're on Windows and have poppler not in PATH, set this
    poppler_bin_path = None
    # e.g. poppler_bin_path = r"C:\poppler-xx\bin"

    pdf_to_images(pdf_file, out_dir, fmt='PNG', dpi=200, poppler_path=poppler_bin_path)
